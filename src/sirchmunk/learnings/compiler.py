# Copyright (c) ModelScope Contributors. All rights reserved.
"""
Knowledge compiler — orchestrates offline compile of document collections.

Fuses PageIndex (tree indexing) and LLM Wiki (knowledge compilation network)
into a single compile pipeline that produces structured tree indices and
knowledge clusters for downstream search acceleration.
"""

import asyncio
import bisect
import ctypes
import gc
import json
import math
import os
import platform
import random
import re
import hashlib
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional, Set, Tuple, Union

from sirchmunk.learnings.tree_indexer import (
    DocumentTree,
    DocumentTreeIndexer,
)
from sirchmunk.llm.openai_chat import OpenAIChat
from sirchmunk.schema.knowledge import (
    AbstractionLevel,
    EvidenceUnit,
    KnowledgeCluster,
    Lifecycle,
    WeakSemanticEdge,
)
from sirchmunk.storage.knowledge_storage import KnowledgeStorage
from sirchmunk.utils import LogCallback, create_logger
from sirchmunk.utils.document_extractor import DocumentExtractor
from sirchmunk.utils.file_utils import get_fast_hash

# Concurrency cap for LLM-heavy file processing
_DEFAULT_CONCURRENCY = 3

# Similarity threshold for merging into existing clusters during compile
_MERGE_SIMILARITY_THRESHOLD = 0.75

# Max chars for manifest-persisted document summary (used in Phase 2 & catalog)
_MANIFEST_SUMMARY_MAX_LEN = 500

# Preview window for direct LLM summarisation (no tree), ~4K tokens
_SUMMARY_PREVIEW_CHARS = 16_000

# Multi-section sampling for large documents without a tree index
_SUMMARY_SAMPLE_SECTIONS = 3          # Number of sections to sample for large docs
_SUMMARY_SAMPLE_SECTION_CHARS = 5_000  # Chars per sampled section

# Targeted table extraction: max chars per table region
_TARGETED_TABLE_MAX_CHARS = 5000

# Targeted table extraction: only process nodes spanning <= N pages
_TABLE_PAGE_SPAN_LIMIT = 5

# Numeric density threshold – fraction of numeric/symbol chars ($, %, digits,
# parenthesised numbers) relative to total non-whitespace chars.  Pages below
# this threshold are skipped during targeted extraction.
_TABLE_NUMERIC_DENSITY_THRESHOLD = 0.15

# Selective force-OCR: max pages to re-extract with forced OCR per document
_FORCE_OCR_MAX_PAGES = 30

# Incremental manifest flush: persist manifest every N completed files
# to survive interrupted compiles without excessive I/O overhead.
_MANIFEST_FLUSH_INTERVAL = 10

# Page-level extraction: max pages to load into memory per batch.
# Prevents loading all 200-400 pages of a large PDF at once.
_PAGE_SCAN_BATCH_SIZE = 50

# How often to run gc.collect() inside the compile loop (every N files).
_GC_INTERVAL = 5


def _force_gc() -> None:
    """Aggressively reclaim Python-managed memory and nudge the C allocator."""
    gc.collect()
    if platform.system() == "Linux":
        try:
            ctypes.CDLL("libc.so.6").malloc_trim(0)
        except (OSError, AttributeError):
            pass


# Shared numeric-token regex for table detection heuristics.
# Matches: $1,234  (1,234)  12.5%  3.14e-5  1,000
_NUM_TOKEN_RE = re.compile(
    r"(?:"
    r"[\$€£¥]\s*[\d,.]+|"
    r"\([\d,.]+\)|"
    r"[\d,.]+%|"
    r"[\d]+\.[\d]+(?:[eE][+-]?\d+)?|"
    r"[\d,]{2,}"
    r")"
)

# A single line with >= this many numeric tokens is treated as a dense
# table row (or multiple rows concatenated), enabling detection even when
# pypdf flattens the entire page to one or two lines.
_DENSE_LINE_MIN_TOKENS = 15

# ---------------------------------------------------------------------------
# Heading normalisation: candidate extraction patterns
# ---------------------------------------------------------------------------
# kreuzberg sometimes renders section titles as ``**bold text**`` or bare
# short standalone lines instead of ``## heading``.  The tree indexer can
# only split on markdown headings, so these "invisible" titles get absorbed
# into parent nodes.
#
# We extract *candidates* via lightweight regexes and let the LLM classify
# which ones are genuine section headings (language/domain-agnostic).

_BOLD_LINE_RE = re.compile(
    r"^\*\*((?:(?!\*\*).)+)\*\*\s*$",
    re.MULTILINE,
)

_STANDALONE_LINE_RE = re.compile(
    r"(?:^|\n\n)([^\n]{5,100})\n\n",
)

_HEADING_CANDIDATE_CAP = 40

# Excel table-level adaptive sampling constants
_XLSX_TOTAL_ROW_BUDGET = 100       # Total sampled rows budget across all sheets
_XLSX_MIN_ROWS_PER_SHEET = 3       # Minimum sampled rows per sheet
_XLSX_MAX_ROWS_PER_SHEET = 50      # Maximum sampled rows per sheet
_XLSX_MAX_SHEETS = 10              # Maximum number of sheets to process
_XLSX_MAX_COLS_DISPLAY = 20        # Maximum columns to display per sheet


# ---------------------------------------------------------------------------
# Data structures
# ---------------------------------------------------------------------------

@dataclass
class FileManifestEntry:
    """State of a single file in the compile manifest."""

    file_hash: str
    compiled_at: str
    has_tree: bool
    cluster_ids: List[str]
    size_bytes: int
    summary: str = ""  # 新增：存储编译期生成的文档摘要
    has_explicit_toc: bool = False  # Whether a native TOC was extracted from the file
    tree_node_count: int = 0  # Number of nodes in the tree index (quality metric)
    has_xlsx_digest: bool = False  # Whether a pre-compiled Excel evidence digest exists
    has_table_digest: bool = False  # Whether PDF tables were extracted and stored
    table_count: int = 0  # Number of tables in this file

    def to_dict(self) -> Dict[str, Any]:
        return {
            "file_hash": self.file_hash,
            "compiled_at": self.compiled_at,
            "has_tree": self.has_tree,
            "cluster_ids": self.cluster_ids,
            "size_bytes": self.size_bytes,
            "summary": self.summary,
            "has_explicit_toc": self.has_explicit_toc,
            "tree_node_count": self.tree_node_count,
            "has_xlsx_digest": self.has_xlsx_digest,
            "has_table_digest": self.has_table_digest,
            "table_count": self.table_count,
        }

    @classmethod
    def from_dict(cls, data: Dict[str, Any]) -> "FileManifestEntry":
        return cls(
            file_hash=data["file_hash"],
            compiled_at=data["compiled_at"],
            has_tree=data.get("has_tree", False),
            cluster_ids=data.get("cluster_ids", []),
            size_bytes=data.get("size_bytes", 0),
            summary=data.get("summary", ""),
            has_explicit_toc=data.get("has_explicit_toc", False),
            tree_node_count=data.get("tree_node_count", 0),
            has_xlsx_digest=data.get("has_xlsx_digest", False),
            has_table_digest=data.get("has_table_digest", False),
            table_count=data.get("table_count", 0),
        )


@dataclass
class CompileManifest:
    """Tracks compiled file states for incremental processing."""

    version: str = "1.0"
    last_compile_at: Optional[str] = None
    files: Dict[str, FileManifestEntry] = field(default_factory=dict)

    def to_json(self) -> str:
        return json.dumps({
            "version": self.version,
            "last_compile_at": self.last_compile_at,
            "files": {k: v.to_dict() for k, v in self.files.items()},
        }, ensure_ascii=False, indent=2)

    @classmethod
    def from_json(cls, json_str: str) -> "CompileManifest":
        data = json.loads(json_str)
        files = {
            k: FileManifestEntry.from_dict(v)
            for k, v in data.get("files", {}).items()
        }
        return cls(
            version=data.get("version", "1.0"),
            last_compile_at=data.get("last_compile_at"),
            files=files,
        )


@dataclass
class FileEntry:
    """Discovered file pending compilation."""

    path: str
    size_bytes: int
    file_hash: str


@dataclass
class ChangeSet:
    """Delta between discovered files and the manifest."""

    added: List[FileEntry] = field(default_factory=list)
    modified: List[FileEntry] = field(default_factory=list)
    deleted: List[str] = field(default_factory=list)
    unchanged: List[str] = field(default_factory=list)


@dataclass
class FileCompileResult:
    """Result of compiling a single file."""

    path: str
    tree: Optional[DocumentTree] = None
    summary: str = ""
    topics: List[str] = field(default_factory=list)
    evidence: Optional[EvidenceUnit] = None
    cluster_ids: List[str] = field(default_factory=list)
    error: Optional[str] = None
    has_explicit_toc: bool = False  # Whether TOC was extracted from native structure
    tree_node_count: int = 0  # Number of nodes in the tree index
    has_xlsx_digest: bool = False  # Whether a pre-compiled Excel evidence digest exists
    has_table_digest: bool = False  # Whether a pre-compiled table digest exists
    table_count: int = 0  # Number of tables extracted


@dataclass
class CompileReport:
    """Summary report of a compile run."""

    total_files: int = 0
    files_added: int = 0
    files_modified: int = 0
    files_skipped: int = 0
    files_deleted: int = 0
    files_sampled: int = 0
    trees_built: int = 0
    clusters_created: int = 0
    clusters_merged: int = 0
    cross_refs_built: int = 0
    errors: List[str] = field(default_factory=list)
    elapsed_seconds: float = 0.0

    def to_dict(self) -> Dict[str, Any]:
        return {
            "total_files": self.total_files,
            "files_added": self.files_added,
            "files_modified": self.files_modified,
            "files_skipped": self.files_skipped,
            "files_deleted": self.files_deleted,
            "files_sampled": self.files_sampled,
            "trees_built": self.trees_built,
            "clusters_created": self.clusters_created,
            "clusters_merged": self.clusters_merged,
            "cross_refs_built": self.cross_refs_built,
            "errors": self.errors,
            "elapsed_seconds": round(self.elapsed_seconds, 2),
        }


@dataclass
class CompileStatus:
    """Status snapshot of the compile state."""

    total_compiled_files: int = 0
    total_clusters: int = 0
    total_trees: int = 0
    last_compile_at: Optional[str] = None
    manifest_path: str = ""


# ---------------------------------------------------------------------------
# Importance probability sampler
# ---------------------------------------------------------------------------

class ImportanceSampler:
    """Select a representative subset of files using importance-based probability.

    Sampling strategy for large datasets:
    - Larger files get higher probability (they contain more information).
    - Uncompiled (new) files are prioritised over previously compiled ones.
    - Files with rare extensions get a mild boost (diversity signal).
    - The final probability is proportional to a composite importance score.
    """

    def __init__(self, max_files: int, seed: Optional[int] = None):
        self._max_files = max_files
        self._rng = random.Random(seed)

    def sample(self, files: List[FileEntry], manifest: CompileManifest) -> List[FileEntry]:
        """Return up to *max_files* entries sampled by importance."""
        if len(files) <= self._max_files:
            return files

        scores = [self._score(f, manifest) for f in files]
        total = sum(scores) or 1.0
        probs = [s / total for s in scores]

        selected_indices = set()
        attempts = 0
        while len(selected_indices) < self._max_files and attempts < len(files) * 3:
            idx = self._weighted_choice(probs)
            selected_indices.add(idx)
            attempts += 1

        return [files[i] for i in sorted(selected_indices)]

    def _score(self, entry: FileEntry, manifest: CompileManifest) -> float:
        """Compute composite importance score."""
        # Size factor: log-scaled, bounded
        size_score = math.log2(max(entry.size_bytes, 1024)) / 20.0

        # Novelty factor: new files are more important
        novelty = 2.0 if entry.path not in manifest.files else 0.5

        # Extension diversity: rare extensions get a mild boost
        ext = Path(entry.path).suffix.lower()
        diversity = 1.5 if ext in {".pdf", ".docx", ".doc", ".tex"} else 1.0

        return size_score * novelty * diversity

    def _weighted_choice(self, probs: List[float]) -> int:
        r = self._rng.random()
        cumulative = 0.0
        for i, p in enumerate(probs):
            cumulative += p
            if r <= cumulative:
                return i
        return len(probs) - 1


# ---------------------------------------------------------------------------
# Compiler
# ---------------------------------------------------------------------------

class KnowledgeCompiler:
    """Orchestrate compile pipeline: file discovery -> tree indexing -> knowledge aggregation."""

    # File extensions eligible for compilation
    _ELIGIBLE_EXTENSIONS = {
        ".pdf", ".docx", ".doc", ".md", ".markdown", ".html", ".htm",
        ".rst", ".tex", ".txt", ".pptx", ".xlsx",
    }

    def __init__(
        self,
        llm: OpenAIChat,
        embedding_client: Optional[Any],
        knowledge_storage: KnowledgeStorage,
        tree_indexer: DocumentTreeIndexer,
        work_path: Union[str, Path],
        log_callback: LogCallback = None,
    ):
        self._llm = llm
        self._embedding = embedding_client
        self._storage = knowledge_storage
        self._tree_indexer = tree_indexer
        self._work_path = Path(work_path).expanduser().resolve()
        self._log = create_logger(log_callback=log_callback)

        self._compile_dir = self._work_path / ".cache" / "compile"
        self._compile_dir.mkdir(parents=True, exist_ok=True)
        self._manifest_path = self._compile_dir / "manifest.json"

    # ------------------------------------------------------------------ #
    #  Resource management                                                #
    # ------------------------------------------------------------------ #

    @staticmethod
    def _configure_thread_limits() -> None:
        """Cap PyTorch thread count to reduce per-thread memory allocation.

        Environment variables (OMP_NUM_THREADS, etc.) are set in the CLI
        entry point before libraries are imported.  This method handles the
        PyTorch-specific runtime API that works retroactively.
        """
        cpu_count = os.cpu_count() or 4
        cap = max(1, min(cpu_count // 2, 4))
        try:
            import torch
            torch.set_num_threads(cap)
            torch.set_num_interop_threads(max(1, cap // 2))
        except (ImportError, RuntimeError):
            pass

    # ------------------------------------------------------------------ #
    #  Public API                                                         #
    # ------------------------------------------------------------------ #

    async def compile(
        self,
        paths: List[str],
        *,
        incremental: bool = True,
        shallow: bool = False,
        max_files: Optional[int] = None,
        concurrency: int = _DEFAULT_CONCURRENCY,
    ) -> CompileReport:
        """Execute the unified knowledge compile pipeline.

        Args:
            paths: Directories or files to compile.
            incremental: Skip unchanged files.
            shallow: Skip tree building even for eligible files — use direct
                     LLM summarisation only (faster, lower quality).
            max_files: Cap on files to process (triggers importance sampling).
            concurrency: Max parallel file compilations.
        """
        import time

        self._configure_thread_limits()

        t0 = time.monotonic()
        report = CompileReport()

        # Phase 1: discover and diff
        await self._log.info("[Compile] Phase 1: File discovery & change detection")
        manifest = self._load_manifest()
        discovered = await self._discover_files(paths)
        report.total_files = len(discovered)
        await self._log.info(f"[Compile] Discovered {len(discovered)} eligible files")

        if incremental:
            changes = self._detect_changes(discovered, manifest)
            to_compile = changes.added + changes.modified
            report.files_skipped = len(changes.unchanged)
            report.files_deleted = len(changes.deleted)

            stale_paths = changes.deleted + [e.path for e in changes.modified]
            if stale_paths:
                await self._purge_stale_artifacts(stale_paths, manifest)
        else:
            to_compile = discovered
            report.files_skipped = 0

        report.files_added = len([f for f in to_compile if f.path not in manifest.files])
        report.files_modified = len(to_compile) - report.files_added

        # Phase 1.5: importance sampling for large datasets
        if max_files and len(to_compile) > max_files:
            await self._log.info(
                f"[Compile] Applying importance sampling: {len(to_compile)} -> {max_files} files"
            )
            sampler = ImportanceSampler(max_files=max_files)
            to_compile = sampler.sample(to_compile, manifest)
            report.files_sampled = len(to_compile)

        if not to_compile:
            await self._log.info("[Compile] No files to compile (all up-to-date)")
            report.elapsed_seconds = time.monotonic() - t0
            return report

        await self._log.info(
            f"[Compile] Phase 2: Processing {len(to_compile)} files "
            f"(concurrency={concurrency})"
        )

        # Phase 2 + 3 (fused): compile files, aggregate inline, release heavy objects
        # Fusing Phase 3 into the completion loop avoids retaining all
        # DocumentTree / EvidenceUnit objects until the end of the pipeline.
        semaphore = asyncio.Semaphore(concurrency)
        _xref_pairs: List[Tuple[str, List[str]]] = []  # lightweight (path, cluster_ids) for Phase 4
        _files_since_flush = 0
        _files_since_gc = 0

        async def _bounded(entry: FileEntry) -> FileCompileResult:
            async with semaphore:
                return await self._compile_single_file(entry, shallow=shallow)

        tasks = [_bounded(f) for f in to_compile]
        for coro in asyncio.as_completed(tasks):
            result = await coro
            if result.error:
                report.errors.append(f"{result.path}: {result.error}")
            else:
                if result.tree:
                    report.trees_built += 1
                manifest.files[result.path] = FileManifestEntry(
                    file_hash=get_fast_hash(result.path) or "",
                    compiled_at=datetime.now(timezone.utc).isoformat(),
                    has_tree=result.tree is not None,
                    cluster_ids=result.cluster_ids,
                    size_bytes=Path(result.path).stat().st_size if Path(result.path).exists() else 0,
                    summary=result.summary[:_MANIFEST_SUMMARY_MAX_LEN] if result.summary else "",
                    has_explicit_toc=result.has_explicit_toc,
                    tree_node_count=result.tree_node_count,
                    has_xlsx_digest=result.has_xlsx_digest,
                    has_table_digest=result.has_table_digest,
                    table_count=result.table_count,
                )

            # Phase 3 inline: aggregate while the result is still alive
            if not result.error and result.summary:
                created, merged = await self._aggregate_to_knowledge_network(result)
                report.clusters_created += created
                report.clusters_merged += merged

            # Retain only lightweight cross-ref data, then drop the heavy result
            _xref_pairs.append((result.path, list(result.cluster_ids)))
            del result

            # Incremental manifest flush to survive interrupted compiles
            _files_since_flush += 1
            if _files_since_flush >= _MANIFEST_FLUSH_INTERVAL:
                manifest.last_compile_at = datetime.now(timezone.utc).isoformat()
                self._save_manifest(manifest)
                _files_since_flush = 0

            _files_since_gc += 1
            if _files_since_gc >= _GC_INTERVAL:
                _force_gc()
                _files_since_gc = 0

        # Phase 2 checkpoint: persist manifest before cross-references
        manifest.last_compile_at = datetime.now(timezone.utc).isoformat()
        self._save_manifest(manifest)

        # Phase 4: cross-references (uses only lightweight path+cluster_ids pairs)
        await self._log.info("[Compile] Phase 4: Building cross-references")
        report.cross_refs_built = await self._build_cross_references_from_pairs(
            _xref_pairs, manifest,
        )

        # Phase 5: persist final manifest + derived indices
        # Catalog and summary index are rebuilt from the manifest, so even
        # partial compiles produce usable search-time metadata.
        manifest.last_compile_at = datetime.now(timezone.utc).isoformat()
        self._save_manifest(manifest)
        self._storage.force_sync()

        self._build_document_catalog(manifest)

        await self._build_summary_index(manifest)

        report.elapsed_seconds = time.monotonic() - t0
        await self._log.info(
            f"[Compile] Done in {report.elapsed_seconds:.1f}s — "
            f"trees={report.trees_built}, created={report.clusters_created}, "
            f"merged={report.clusters_merged}, errors={len(report.errors)}"
        )
        return report

    async def get_status(self, paths: List[str]) -> CompileStatus:
        """Return current compile status for the given paths."""
        manifest = self._load_manifest()
        path_set = {str(Path(p).resolve()) for p in paths}

        compiled_count = 0
        tree_count = 0
        cluster_ids: Set[str] = set()
        for fp, entry in manifest.files.items():
            for p in path_set:
                if fp.startswith(p):
                    compiled_count += 1
                    if entry.has_tree:
                        tree_count += 1
                    cluster_ids.update(entry.cluster_ids)
                    break

        return CompileStatus(
            total_compiled_files=compiled_count,
            total_clusters=len(cluster_ids),
            total_trees=tree_count,
            last_compile_at=manifest.last_compile_at,
            manifest_path=str(self._manifest_path),
        )

    # ------------------------------------------------------------------ #
    #  File discovery and change detection                                #
    # ------------------------------------------------------------------ #

    async def _discover_files(self, paths: List[str]) -> List[FileEntry]:
        """Walk paths and return all compilation-eligible files."""
        entries: List[FileEntry] = []
        seen: Set[str] = set()

        for base in paths:
            base_path = Path(base).expanduser().resolve()
            if base_path.is_file():
                candidates = [base_path]
            elif base_path.is_dir():
                candidates = sorted(base_path.rglob("*"))
            else:
                continue

            for fp in candidates:
                if not fp.is_file():
                    continue
                if fp.suffix.lower() not in self._ELIGIBLE_EXTENSIONS:
                    continue
                abs_path = str(fp.resolve())
                if abs_path in seen:
                    continue
                seen.add(abs_path)
                fh = get_fast_hash(abs_path)
                if fh is None:
                    continue
                entries.append(FileEntry(
                    path=abs_path,
                    size_bytes=fp.stat().st_size,
                    file_hash=fh,
                ))

        return entries

    def _detect_changes(
        self, discovered: List[FileEntry], manifest: CompileManifest,
    ) -> ChangeSet:
        """Compare discovered files against the manifest for incremental compile."""
        changes = ChangeSet()
        current_paths = {f.path for f in discovered}

        for entry in discovered:
            prev = manifest.files.get(entry.path)
            if prev is None:
                changes.added.append(entry)
            elif prev.file_hash != entry.file_hash:
                changes.modified.append(entry)
            else:
                changes.unchanged.append(entry.path)

        for old_path in manifest.files:
            if old_path not in current_paths:
                changes.deleted.append(old_path)

        return changes

    # ------------------------------------------------------------------ #
    #  Stale artifact cleanup                                             #
    # ------------------------------------------------------------------ #

    async def _purge_stale_artifacts(
        self,
        file_paths: List[str],
        manifest: CompileManifest,
    ) -> None:
        """Remove disk artifacts and DuckDB clusters for deleted/modified files.

        Called before recompilation so that modified files start with a
        clean slate and deleted files leave no residue.
        """
        artifact_dirs = {
            "trees": ".json",
            "content": ".txt",
            "table_digests": ".json",
            "xlsx_digests": ".txt",
        }

        for file_path in file_paths:
            entry = manifest.files.get(file_path)
            if entry is None:
                continue

            file_hash = entry.file_hash

            # 1. Remove disk artifacts keyed by file_hash
            if file_hash:
                for subdir, ext in artifact_dirs.items():
                    artifact = self._compile_dir / subdir / f"{file_hash}{ext}"
                    try:
                        artifact.unlink(missing_ok=True)
                    except OSError:
                        pass

            # 2. Remove associated knowledge clusters from DuckDB
            for cluster_id in entry.cluster_ids:
                try:
                    await self._storage.remove(cluster_id)
                except Exception:
                    pass

            # 3. Drop the manifest entry
            manifest.files.pop(file_path, None)

    # ------------------------------------------------------------------ #
    #  Single-file compilation                                            #
    # ------------------------------------------------------------------ #

    async def _compile_single_file(
        self,
        entry: FileEntry,
        *,
        shallow: bool = False,
    ) -> FileCompileResult:
        """Unified compile pipeline: tree-if-eligible -> summary -> topics -> evidence.

        When *shallow* is True (or file is ineligible for tree indexing),
        the pipeline skips tree building and summarises via a direct LLM call.

        Large intermediate objects (extraction output, enriched content,
        raw tables) are explicitly released after their last use to keep
        per-file peak memory bounded.
        """
        result = FileCompileResult(path=entry.path)
        try:
            await self._log.info(f"[Compile] Processing: {Path(entry.path).name}")

            extraction = await DocumentExtractor.extract_isolated(
                entry.path, DocumentExtractor.ENHANCED,
            )
            content = extraction.content
            content = await self._normalize_bold_headings(content)
            if not content or len(content.strip()) < 100:
                result.error = "Insufficient text content"
                return result

            # Extract scalar metadata from extraction before releasing it
            page_count = extraction.page_count
            raw_tables = extraction.tables
            del extraction

            use_tree = (
                not shallow
                and DocumentTreeIndexer.should_build_tree(entry.path, len(content))
            )

            # Phase 0.5: TOC extraction (layers 1-3 are zero LLM calls)
            toc_entries = None
            if use_tree:
                from sirchmunk.learnings.toc_extractor import TOCExtractor
                toc_entries = await TOCExtractor.extract(
                    entry.path, content,
                    total_pages=page_count,
                )
                if toc_entries:
                    await self._log.info(
                        f"[Compile] Extracted TOC with {len(toc_entries)} entries "
                        f"for {Path(entry.path).name}"
                    )

            if use_tree:
                result.tree = await self._tree_indexer.build_tree(
                    entry.path, content,
                    toc_entries=toc_entries,
                    total_pages=page_count,
                )

            result.has_explicit_toc = bool(toc_entries)
            del toc_entries
            result.tree_node_count = self._count_tree_nodes(result.tree)
            print(f"SEARCH_WIKI_DEBUG [C2] tree_build: success={result.tree is not None}, nodes={result.tree_node_count}, tree.file_path={result.tree.file_path if result.tree else 'N/A'}", flush=True)

            # --- Summary + topics + evidence (needs content) ---
            ext = Path(entry.path).suffix.lower()
            evidence_digest = ""

            if ext in (".xlsx", ".xls"):
                metadata_prefix, evidence_digest = self._extract_xlsx_sampling(entry.path)
            else:
                metadata_prefix = self._extract_structured_metadata(entry.path, content)

            # Build enriched_content only for the summary LLM call, then release
            if metadata_prefix:
                result.summary = await self._extract_summary(
                    entry.path, metadata_prefix + content, result.tree,
                )
            else:
                result.summary = await self._extract_summary(
                    entry.path, content, result.tree,
                )
            del metadata_prefix

            result.topics = await self._extract_topics(result.summary)
            result.evidence = self._build_evidence(entry, content, result)

            # Persist Excel evidence digest
            if evidence_digest.strip():
                try:
                    digest_dir = self._compile_dir / "xlsx_digests"
                    digest_dir.mkdir(parents=True, exist_ok=True)
                    file_hash = get_fast_hash(entry.path) or ""
                    if file_hash:
                        (digest_dir / f"{file_hash}.txt").write_text(
                            evidence_digest, encoding="utf-8",
                        )
                        result.has_xlsx_digest = True
                except Exception:
                    pass
            del evidence_digest

            # Cache ENHANCED content to disk
            try:
                file_hash_content = get_fast_hash(entry.path) or ""
                if file_hash_content and content:
                    content_dir = self._compile_dir / "content"
                    content_dir.mkdir(parents=True, exist_ok=True)
                    (content_dir / f"{file_hash_content}.txt").write_text(
                        content, encoding="utf-8",
                    )
            except Exception:
                pass

            # --- Table digest + integration (needs raw_tables, then release) ---
            if raw_tables:
                try:
                    table_digest = self._build_table_digest(raw_tables)
                    if table_digest:
                        digest_dir = self._compile_dir / "table_digests"
                        digest_dir.mkdir(parents=True, exist_ok=True)
                        file_hash = get_fast_hash(entry.path) or ""
                        if file_hash:
                            (digest_dir / f"{file_hash}.json").write_text(
                                json.dumps(table_digest, ensure_ascii=False),
                                encoding="utf-8",
                            )
                            result.has_table_digest = True
                            result.table_count = len(raw_tables)
                except Exception:
                    pass

                if result.tree and result.tree.root:
                    self._integrate_tables_into_tree(
                        result.tree.root, raw_tables,
                        content=content, total_pages=page_count,
                    )

            print(f"SEARCH_WIKI_DEBUG [C3] table_digest: generated={result.has_table_digest}, count={result.table_count}", flush=True)
            del raw_tables

            # --- Phases 2.5-2.8: secondary table extraction (PDF only) ---
            # These phases re-read from the PDF file; `content` is only
            # needed for Phase 2.6 fallback and Phase 2.8 enrichment.
            if ext == ".pdf":
                if result.tree and result.tree.root:
                    targeted_tables = await self._targeted_table_extraction(
                        entry.path, result.tree,
                    )
                    await self._supplement_table_digest(
                        entry.path, targeted_tables, result,
                        source_label="Targeted extraction",
                    )
                    del targeted_tables

                if page_count:
                    covered_pages = self._get_covered_table_pages(entry.path)
                    tree_root = (
                        result.tree.root
                        if result.tree and result.tree.root else None
                    )
                    content_tables = await self._content_based_table_scan(
                        entry.path, page_count, covered_pages,
                        enhanced_content=content, tree_root=tree_root,
                    )
                    await self._supplement_table_digest(
                        entry.path, content_tables, result,
                        source_label="Content-based scan",
                    )
                    del content_tables

                    covered_after_scan = self._get_covered_table_pages(entry.path)
                    gap_pages = self._find_force_ocr_candidates(
                        entry.path, page_count, covered_after_scan,
                    )
                    if gap_pages:
                        ocr_tables = await self._selective_force_ocr_tables(
                            entry.path, gap_pages,
                        )
                        await self._supplement_table_digest(
                            entry.path, ocr_tables, result,
                            source_label="Selective force-OCR",
                        )
                        del ocr_tables

                if result.has_table_digest:
                    self._enrich_table_digest_content(
                        entry.path, content, tree_root=None,
                    )

            # Content is no longer needed — release before returning
            del content

        except Exception as exc:
            result.error = str(exc)
            await self._log.warning(f"[Compile] Failed: {entry.path}: {exc}")

        return result

    @staticmethod
    def _is_generic_summary(summary: str, min_specificity_len: int = 80) -> bool:
        """Check whether a summary is too generic to be useful for retrieval.

        A generic summary typically contains only structural descriptions
        (e.g., "This document contains several sections") without specific
        content indicators.  Detection uses summary length and information
        density as domain-agnostic proxies.
        """
        if not summary:
            return True
        stripped = summary.strip()
        if len(stripped) < min_specificity_len:
            return True
        # Count unique substantive words (>4 chars) as a proxy for specificity
        words = set(w.lower() for w in stripped.split() if len(w) > 4)
        return len(words) < 8

    async def _extract_summary(
        self,
        file_path: str,
        content: str,
        tree: Optional[DocumentTree] = None,
    ) -> str:
        """Generate a document-level summary.

        When a tree is available its root already contains an LLM-synthesized
        summary (produced by ``_synthesize_root_summary`` during tree build),
        so we reuse it directly — unless the summary is too generic (Plan 2),
        in which case we fall back to multi-section LLM summarization.

        For large documents without a tree, uses multi-section sampling
        (beginning, middle, end) to capture the full scope of the document.
        """
        if tree and tree.root and tree.root.summary:
            if not self._is_generic_summary(tree.root.summary):
                return tree.root.summary
            await self._log.info(
                f"[Compile] Root summary too generic for {Path(file_path).name}, "
                f"falling back to LLM summarization"
            )

        preview = self._build_summary_preview(content)
        from sirchmunk.llm.prompts import COMPILE_DOC_SUMMARY
        prompt = COMPILE_DOC_SUMMARY.format(
            file_name=Path(file_path).name,
            document_content=preview,
        )
        resp = await self._llm.achat([{"role": "user", "content": prompt}])
        return resp.content.strip()

    @staticmethod
    def _build_summary_preview(content: str) -> str:
        """Build a representative preview for LLM summarisation.

        For short documents (≤ _SUMMARY_PREVIEW_CHARS), returns the full
        content.  For large documents, samples the beginning, middle, and
        end to capture the document's full scope within the token budget.
        """
        if len(content) <= _SUMMARY_PREVIEW_CHARS:
            return content

        section_size = _SUMMARY_SAMPLE_SECTION_CHARS
        mid_start = max(section_size, (len(content) - section_size) // 2)

        head = content[:section_size]
        middle = content[mid_start:mid_start + section_size]
        tail = content[-section_size:]

        return (
            f"[Beginning of document]\n{head}\n\n"
            f"[... content omitted ...]\n\n"
            f"[Middle of document]\n{middle}\n\n"
            f"[... content omitted ...]\n\n"
            f"[End of document]\n{tail}"
        )

    @staticmethod
    def _extract_structured_metadata(file_path: str, content: str) -> str:
        """Extract structural metadata for non-text document types.

        For spreadsheets and presentations, prepend a structural overview
        (sheet names, column headers, slide titles) so the LLM summariser
        has better context than raw extracted text alone.

        Returns a metadata prefix string (may be empty for unsupported types).
        """
        ext = Path(file_path).suffix.lower()

        if ext == ".xlsx":
            metadata, _evidence = KnowledgeCompiler._extract_xlsx_sampling(file_path)
            return metadata
        if ext == ".pptx":
            return KnowledgeCompiler._extract_pptx_metadata(file_path)

        return ""

    @staticmethod
    def _compute_xlsx_sample_rows(total_rows: int, num_sheets: int, sheet_rows: int) -> int:
        """Compute adaptive sample row count per sheet.

        Strategy:
        - Divides _XLSX_TOTAL_ROW_BUDGET equally across sheets
        - Small sheets (<=budget) are fully sampled
        - Large sheets are capped at budget
        - Result clamped to [_XLSX_MIN_ROWS_PER_SHEET, _XLSX_MAX_ROWS_PER_SHEET]
        """
        budget_per_sheet = max(1, _XLSX_TOTAL_ROW_BUDGET // max(1, num_sheets))
        n = min(sheet_rows, budget_per_sheet)
        return max(_XLSX_MIN_ROWS_PER_SHEET, min(_XLSX_MAX_ROWS_PER_SHEET, n))

    @staticmethod
    def _extract_xlsx_sampling(file_path: str) -> Tuple[str, str]:
        """Extract structural metadata AND sampled content from Excel workbook.

        Performs table-level intelligent sampling with adaptive row counts
        based on workbook size and sheet complexity.

        Returns:
            (metadata_prefix, evidence_digest)
            - metadata_prefix: injected into summary generation context
            - evidence_digest: structured text usable directly as search evidence
        """
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)

            sheet_names = wb.sheetnames[:_XLSX_MAX_SHEETS]
            num_sheets = len(sheet_names)

            # Phase 1: Collect sheet statistics
            sheet_stats: List[Dict[str, Any]] = []
            for sheet_name in sheet_names:
                ws = wb[sheet_name]
                row_count = ws.max_row or 0
                col_count = ws.max_column or 0
                # Read headers (first row)
                headers: List[str] = []
                for row in ws.iter_rows(min_row=1, max_row=1, values_only=True):
                    headers = [str(h) for h in row if h is not None]
                    break
                sheet_stats.append({
                    "name": sheet_name,
                    "rows": row_count,
                    "cols": col_count,
                    "headers": headers[:_XLSX_MAX_COLS_DISPLAY],
                    "ws": ws,
                })

            # Phase 2: Calculate total rows for adaptive sampling
            total_rows = sum(s["rows"] for s in sheet_stats)

            meta_lines: List[str] = ["[Excel Workbook Structure]"]
            evidence_lines: List[str] = []

            for stat in sheet_stats:
                ws = stat["ws"]
                sheet_name = stat["name"]
                row_count = stat["rows"]
                col_count = stat["cols"]
                headers = stat["headers"]
                header_str = ", ".join(headers) if headers else "no headers"

                # Metadata line
                meta_lines.append(
                    f"- Sheet '{sheet_name}': {row_count} rows, {col_count} columns, "
                    f"headers: [{header_str}]"
                )

                # Adaptive sampling
                sample_n = KnowledgeCompiler._compute_xlsx_sample_rows(
                    total_rows, num_sheets, row_count
                )

                evidence_lines.append(
                    f"[Sheet '{sheet_name}' ({row_count} rows, {col_count} columns)]"
                )
                evidence_lines.append(f"Columns: {header_str}")

                # Sample rows
                if row_count <= sample_n:
                    evidence_lines.append(f"(Full content - {row_count} rows)")
                else:
                    evidence_lines.append(f"Sample rows (top {sample_n} of {row_count}):")

                # Build table header
                display_headers = headers[:_XLSX_MAX_COLS_DISPLAY]
                if display_headers:
                    evidence_lines.append("| " + " | ".join(display_headers) + " |")
                    evidence_lines.append("|" + "|".join(["---"] * len(display_headers)) + "|")

                # Read sample rows (skip header row)
                numeric_cols: Dict[int, List[float]] = {}  # col_index -> numeric values
                sampled = 0
                for row in ws.iter_rows(
                    min_row=2,
                    max_row=min(row_count, sample_n + 1),
                    values_only=True,
                ):
                    cells: List[str] = []
                    for ci, cell_val in enumerate(row):
                        if ci >= _XLSX_MAX_COLS_DISPLAY:
                            break
                        str_val = str(cell_val) if cell_val is not None else ""
                        cells.append(str_val[:50])  # truncate long cell values
                        # Track numeric values for statistics
                        if isinstance(cell_val, (int, float)) and cell_val == cell_val:
                            numeric_cols.setdefault(ci, []).append(float(cell_val))
                    if cells:
                        evidence_lines.append("| " + " | ".join(cells) + " |")
                    sampled += 1

                # Statistics for numeric columns
                stat_parts: List[str] = []
                for ci, values in numeric_cols.items():
                    if len(values) >= 2 and ci < len(display_headers):
                        col_name = display_headers[ci]
                        stat_parts.append(
                            f"{col_name} range [{min(values):.4g}-{max(values):.4g}]"
                        )
                if stat_parts:
                    evidence_lines.append(f"Statistics: {', '.join(stat_parts[:5])}")

                evidence_lines.append("")  # blank line between sheets

            wb.close()

            metadata = "\n".join(meta_lines) + "\n\n"
            evidence = "\n".join(evidence_lines)
            return metadata, evidence

        except Exception:
            return "", ""

    @staticmethod
    def _extract_xlsx_metadata(file_path: str) -> str:
        """Extract structural metadata from Excel files (legacy wrapper).

        Delegates to _extract_xlsx_sampling and returns only the metadata prefix
        for backward compatibility.
        """
        metadata, _evidence = KnowledgeCompiler._extract_xlsx_sampling(file_path)
        return metadata

    @staticmethod
    def _extract_pptx_metadata(file_path: str) -> str:
        """Extract structural metadata from PowerPoint files.

        Reads slide count and titles (from the title placeholder) to give
        the LLM a table-of-contents-like overview of the presentation.
        Caps at 20 slides for bounded output.
        """
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            lines: List[str] = [f"[PowerPoint Structure: {len(prs.slides)} slides]"]
            for i, slide in enumerate(prs.slides[:20], 1):  # Cap at 20 slides
                title = ""
                if slide.shapes.title:
                    title = slide.shapes.title.text.strip()
                if title:
                    lines.append(f"- Slide {i}: {title}")
            return "\n".join(lines) + "\n\n"
        except Exception:
            return ""

    def _build_evidence(
        self,
        entry: FileEntry,
        content: str,
        result: FileCompileResult,
    ) -> EvidenceUnit:
        """Build an EvidenceUnit, populating snippets/tree_path from tree leaves."""
        from sirchmunk.schema.metadata import FileInfo

        snippets: List[str] = []
        tree_path: Optional[List[str]] = None

        if result.tree and result.tree.root:
            leaves = result.tree.root.all_leaves()
            tree_path = [leaf.node_id for leaf in leaves]
            for leaf in leaves:
                start, end = leaf.char_range
                snippet = content[start:end][:500]
                if snippet.strip():
                    snippets.append(snippet)

        return EvidenceUnit(
            doc_id=FileInfo.get_cache_key(entry.path),
            file_or_url=Path(entry.path),
            summary=result.summary,
            is_found=True,
            snippets=snippets,
            tree_path=tree_path,
            extracted_at=datetime.now(timezone.utc),
        )

    async def _extract_topics(self, summary: str) -> List[str]:
        """Extract key topics/entities from a document summary."""
        from sirchmunk.llm.prompts import COMPILE_TOPIC_EXTRACTION
        prompt = COMPILE_TOPIC_EXTRACTION.format(summary=summary)
        resp = await self._llm.achat([{"role": "user", "content": prompt}])
        try:
            raw = resp.content.strip()
            if raw.startswith("["):
                parsed = json.loads(raw)
                if isinstance(parsed, list):
                    return [str(t) for t in parsed if t]
            return [t.strip() for t in raw.split(",") if t.strip()]
        except (json.JSONDecodeError, TypeError):
            return []

    # ------------------------------------------------------------------ #
    #  Knowledge aggregation (LLM Wiki Ingest)                            #
    # ------------------------------------------------------------------ #

    async def _aggregate_to_knowledge_network(
        self, result: FileCompileResult,
    ) -> Tuple[int, int]:
        """Aggregate a file's compile result into the knowledge network.

        Three-tier similarity strategy (per design doc):
          - similarity >= 0.80  → merge into existing cluster
          - 0.50 <= sim < 0.80  → create new cluster + weak edge to similar
          - similarity < 0.50   → create standalone cluster

        Returns:
            (clusters_created, clusters_merged)
        """
        created, merged = 0, 0
        if not result.summary:
            return created, merged

        embedding = self._encode_text(result.summary)

        # Search for similar existing clusters across a wider range
        best_match: Optional[Dict[str, Any]] = None
        if embedding is not None:
            similar = await self._storage.search_similar_clusters(
                query_embedding=embedding,
                top_k=3,
                similarity_threshold=0.50,
            )
            if similar:
                best_match = similar[0]

        if best_match and best_match["similarity"] >= 0.80:
            # Tier 1: merge into existing cluster
            cluster = await self._storage.get(best_match["id"])
            if cluster:
                await self._merge_into_cluster(cluster, result)
                # Re-compute embedding for merged content
                await self._update_cluster_embedding(cluster)
                result.cluster_ids.append(cluster.id)
                merged += 1
                return created, merged

        # Create a new cluster (Tier 2 or Tier 3)
        cluster = await self._create_cluster(result)
        if cluster:
            result.cluster_ids.append(cluster.id)
            await self._store_cluster_embedding(cluster, embedding, result.summary)
            created += 1

            # Tier 2: build weak edges to moderately similar clusters
            if best_match and best_match["similarity"] >= 0.50:
                for s in (similar or []):
                    if s["similarity"] >= 0.50:
                        target = await self._storage.get(s["id"])
                        if target:
                            self._add_edge(cluster, target.id, "embed_sim", s["similarity"])
                            self._add_edge(target, cluster.id, "embed_sim", s["similarity"])
                            await self._storage.update(target)
                await self._storage.update(cluster)

        return created, merged

    def _encode_text(self, text: str) -> Optional[Any]:
        """Encode text to embedding vector, returns None on failure."""
        if not self._embedding or not self._embedding.is_ready():
            return None
        try:
            vectors = self._embedding._encode_sync([text])
            return vectors[0] if len(vectors) > 0 else None
        except Exception:
            return None

    async def _store_cluster_embedding(
        self, cluster: KnowledgeCluster, embedding: Optional[Any], text: str,
    ) -> None:
        """Store embedding for a cluster if available."""
        if embedding is None or not self._embedding:
            return
        text_hash = hashlib.md5(text.encode()).hexdigest()
        vec = embedding.tolist() if hasattr(embedding, "tolist") else list(embedding)
        await self._storage.store_embedding(
            cluster.id, vec,
            self._embedding.model_id or "default",
            text_hash,
        )

    async def _update_cluster_embedding(self, cluster: KnowledgeCluster) -> None:
        """Re-compute and store embedding after content merge."""
        content_text = str(cluster.content)[:2000] if cluster.content else ""
        if not content_text:
            return
        embedding = self._encode_text(content_text)
        await self._store_cluster_embedding(cluster, embedding, content_text)

    async def _merge_into_cluster(
        self,
        cluster: KnowledgeCluster,
        result: FileCompileResult,
    ) -> None:
        """Merge a file compile result into an existing cluster."""
        # Append evidence
        if result.evidence:
            existing_doc_ids = {e.doc_id for e in cluster.evidences}
            if result.evidence.doc_id not in existing_doc_ids:
                cluster.evidences.append(result.evidence)

        # Enrich content via LLM merge
        from sirchmunk.llm.prompts import COMPILE_MERGE_KNOWLEDGE
        prompt = COMPILE_MERGE_KNOWLEDGE.format(
            existing_content=str(cluster.content)[:3000],
            new_summary=result.summary[:3000],
        )
        resp = await self._llm.achat([{"role": "user", "content": prompt}])
        cluster.content = resp.content.strip()

        # Update metadata
        cluster.search_results = list(set(
            (cluster.search_results or []) + [result.path]
        ))
        merge_count = getattr(cluster, "merge_count", 0) or 0
        cluster.merge_count = merge_count + 1

        # Lifecycle promotion
        if cluster.merge_count >= 3 and cluster.lifecycle == Lifecycle.EMERGING:
            cluster.lifecycle = Lifecycle.STABLE

        await self._storage.update(cluster)

    async def _create_cluster(
        self, result: FileCompileResult,
    ) -> Optional[KnowledgeCluster]:
        """Create a new KnowledgeCluster from a file compile result."""
        cluster_text = result.summary
        cluster_id = f"C{hashlib.sha256(cluster_text.encode('utf-8')).hexdigest()[:10]}"

        name = Path(result.path).stem[:60]
        if result.topics:
            name = result.topics[0][:60]

        cluster = KnowledgeCluster(
            id=cluster_id,
            name=name,
            description=[result.summary[:500]],
            content=result.summary,
            evidences=[result.evidence] if result.evidence else [],
            patterns=result.topics[:5],
            lifecycle=Lifecycle.EMERGING,
            confidence=0.5,
            abstraction_level=AbstractionLevel.TECHNIQUE,
            hotness=0.3,
            search_results=[result.path],
        )

        ok = await self._storage.insert(cluster)
        return cluster if ok else None

    # ------------------------------------------------------------------ #
    #  Cross-references                                                   #
    # ------------------------------------------------------------------ #

    async def _build_cross_references_from_pairs(
        self,
        pairs: List[Tuple[str, List[str]]],
        manifest: CompileManifest,
    ) -> int:
        """Build co-occurrence edges between clusters that share source files.

        Accepts lightweight ``(path, cluster_ids)`` pairs instead of full
        ``FileCompileResult`` objects to avoid retaining heavy compile results.
        Includes historical data from the manifest.
        """
        cluster_to_files: Dict[str, Set[str]] = {}

        for path, cluster_ids in pairs:
            for cid in cluster_ids:
                cluster_to_files.setdefault(cid, set()).add(path)

        for fp, entry in manifest.files.items():
            for cid in entry.cluster_ids:
                cluster_to_files.setdefault(cid, set()).add(fp)

        # Find cluster pairs that share at least one source file
        cluster_ids = list(cluster_to_files.keys())
        edges_created = 0
        pairs_seen: Set[Tuple[str, str]] = set()

        for i in range(len(cluster_ids)):
            for j in range(i + 1, len(cluster_ids)):
                cid_a, cid_b = cluster_ids[i], cluster_ids[j]
                shared = cluster_to_files[cid_a] & cluster_to_files[cid_b]
                if not shared:
                    continue

                pair_key = (min(cid_a, cid_b), max(cid_a, cid_b))
                if pair_key in pairs_seen:
                    continue
                pairs_seen.add(pair_key)

                weight = min(len(shared) * 0.25, 1.0)
                c_a = await self._storage.get(cid_a)
                c_b = await self._storage.get(cid_b)
                if c_a and c_b:
                    self._add_edge(c_a, cid_b, "co_occur", weight)
                    self._add_edge(c_b, cid_a, "co_occur", weight)
                    await self._storage.update(c_a)
                    await self._storage.update(c_b)
                    edges_created += 1

        return edges_created

    @staticmethod
    def _add_edge(
        cluster: KnowledgeCluster, target_id: str, source: str, weight: float,
    ) -> None:
        """Add or update a WeakSemanticEdge on a cluster."""
        for edge in cluster.related_clusters:
            if edge.target_cluster_id == target_id and edge.source == source:
                edge.weight = max(edge.weight, weight)
                return
        cluster.related_clusters.append(
            WeakSemanticEdge(target_cluster_id=target_id, weight=weight, source=source)
        )

    def _build_table_digest(
        self, tables: List[Dict[str, Any]],
    ) -> Optional[Dict[str, Any]]:
        """Build a structured table digest from extraction output.

        Returns a versioned JSON-serializable dict containing all tables
        with their page numbers, markdown representation, and cell data.
        Tables are indexed for page-range-based retrieval at search time.
        """
        if not tables:
            return None

        digest_tables = []
        for idx, table in enumerate(tables):
            markdown = table.get("markdown", "")
            cells = table.get("cells", [])
            if not markdown and not cells:
                continue

            # Compute row/col counts from cells (kreuzberg returns List[List[str]])
            row_count = 0
            col_count = 0
            if cells:
                row_count = len(cells)
                col_count = max((len(row) for row in cells if isinstance(row, (list, tuple))), default=0)
            elif markdown:
                # Estimate from markdown lines
                lines = [l for l in markdown.strip().split("\n") if l.strip().startswith("|")]
                row_count = max(0, len(lines) - 1)  # exclude separator
                col_count = lines[0].count("|") - 1 if lines else 0

            # Skip pseudo-tables: single-column or insufficient structure
            if col_count <= 1:
                continue

            digest_tables.append({
                "index": idx,
                "page_number": table.get("page_number"),
                "markdown": markdown,
                "row_count": row_count,
                "col_count": col_count,
                "cells": cells,
            })

        if not digest_tables:
            return None

        return {
            "version": 1,
            "table_count": len(digest_tables),
            "tables": digest_tables,
        }

    def _integrate_tables_into_tree(
        self,
        node: "TreeNode",
        tables: List[Dict[str, Any]],
        content: str,
        *,
        total_pages: Optional[int] = None,
        _counter: Optional[List[int]] = None,
    ) -> None:
        """Integrate tables into tree: annotate counts AND create table child nodes for leaf nodes.

        For each node with a valid page_range, counts how many valid extracted
        tables fall within that range (excluding pseudo-tables with col_count <= 1).
        For leaf nodes with matching tables, creates dedicated TreeNode children
        with ``content_type="table"``.
        """
        from sirchmunk.learnings.tree_indexer import TreeNode

        if node is None:
            return

        if _counter is None:
            _counter = [0]

        # Depth-first: process existing children first
        for child in list(node.children):
            self._integrate_tables_into_tree(
                child, tables, content,
                total_pages=total_pages, _counter=_counter,
            )

        # Match valid tables to this node's page_range
        matched_tables: List[Dict[str, Any]] = []
        if node.page_range:
            ps, pe = node.page_range
            for t in tables:
                pn = t.get("page_number")
                if pn is None or not (ps <= pn <= pe):
                    continue
                # Skip pseudo-tables
                if self._is_pseudo_table(t):
                    continue
                matched_tables.append(t)

        node.table_count = len(matched_tables)

        # NOTE: _spawn_table_children disabled - converting leaf to non-leaf breaks
        # search navigation which expects leaves for char_range extraction.
        # TODO: Re-enable when search can properly handle mixed text+table children.
        # if not node.children and matched_tables:
        #     try:
        #         self._spawn_table_children(
        #             node, matched_tables, content, _counter,
        #         )
        #     except Exception:
        #         pass

    @staticmethod
    def _is_pseudo_table(table: Dict[str, Any]) -> bool:
        """Return True if the table lacks meaningful structure (col_count <= 1)."""
        markdown = table.get("markdown", "")
        cells = table.get("cells", [])
        if not markdown and not cells:
            return True
        col_count = 0
        if cells:
            col_count = max(
                (len(row) for row in cells if isinstance(row, (list, tuple))),
                default=0,
            )
        elif markdown:
            lines = [l for l in markdown.strip().split("\n") if l.strip().startswith("|")]
            col_count = (lines[0].count("|") - 1) if lines else 0
        return col_count <= 1

    def _spawn_table_children(
        self,
        node: "TreeNode",
        matched_tables: List[Dict[str, Any]],
        content: str,
        counter: List[int],
    ) -> None:
        """Create TreeNode children for each matched table under a leaf node.

        Also inserts a text-content sibling preserving the original leaf content.
        """
        from sirchmunk.learnings.tree_indexer import TreeNode

        child_level = node.level + 1

        # Preserve original text content as first child
        text_child_id = f"T{counter[0]:06d}"
        counter[0] += 1
        node.children.append(
            TreeNode(
                node_id=text_child_id,
                title=node.title,
                summary=node.summary[:300] if node.summary else "",
                char_range=node.char_range,
                level=child_level,
                page_range=node.page_range,
                children=[],
                table_count=0,
                content_type="text",
            )
        )

        # Create one child per table
        for table in matched_tables:
            tid = f"T{counter[0]:06d}"
            counter[0] += 1

            markdown = table.get("markdown", "")
            title = self._extract_table_title(table)
            page_number = table.get("page_number")

            # Attempt to locate table markdown in content
            char_range = node.char_range
            if markdown and content:
                pos = content.find(markdown[:120])
                if pos >= 0:
                    char_range = (pos, pos + len(markdown))

            page_range = (
                (page_number, page_number) if page_number is not None
                else node.page_range
            )

            node.children.append(
                TreeNode(
                    node_id=tid,
                    title=title,
                    summary=markdown[:300] if markdown else "",
                    char_range=char_range,
                    level=child_level,
                    page_range=page_range,
                    children=[],
                    table_count=0,
                    content_type="table",
                )
            )

    @staticmethod
    def _extract_table_title(table: Dict[str, Any]) -> str:
        """Extract a concise title from table markdown header row.

        Parses the first meaningful line of the markdown table (skipping
        separator rows like ``|---|---|``), strips ``|`` delimiters, and
        returns the first 80 characters as the title.
        """
        markdown = table.get("markdown", "")
        if not markdown:
            pn = table.get("page_number", "?")
            return f"Table (p.{pn})"

        for line in markdown.strip().split("\n"):
            stripped = line.strip()
            if not stripped:
                continue
            # Skip separator rows (e.g. |---|---| or +---+---+)
            content_chars = stripped.replace("|", "").replace("-", "").replace(":", "").replace("+", "").strip()
            if not content_chars:
                continue
            # Extract cell contents
            title = " | ".join(
                seg.strip() for seg in stripped.split("|") if seg.strip()
            )
            return title[:80] if title else f"Table (p.{table.get('page_number', '?')})"

        pn = table.get("page_number", "?")
        return f"Table (p.{pn})"

    @staticmethod
    def _count_tree_nodes(tree: Optional[DocumentTree]) -> int:
        """Count total nodes in a DocumentTree (recursive).

        Args:
            tree: The tree to count, or None.

        Returns:
            Total node count, or 0 if tree is None.
        """
        if tree is None or tree.root is None:
            return 0

        def _count(node: Any) -> int:
            return 1 + sum(_count(c) for c in node.children)

        return _count(tree.root)

    # ------------------------------------------------------------------ #
    #  Targeted table extraction                                          #
    # ------------------------------------------------------------------ #

    async def _targeted_table_extraction(
        self, file_path: str, tree: DocumentTree,
    ) -> list[dict]:
        """Extract tables from tree nodes likely containing tabular data.

        Uses generic structural signals (metadata, page span, numeric
        density) instead of domain-specific title keywords.  For each
        candidate with a valid ``page_range``, extracts per-page text
        via :meth:`DocumentExtractor.extract_page_range` and applies
        heuristic table-region detection.  Pages whose numeric density
        falls below ``_TABLE_NUMERIC_DENSITY_THRESHOLD`` are skipped.

        Returns:
            List of table dicts compatible with the table-digest format::

                {"page": int, "content": str, "source": str}
        """
        if tree is None or tree.root is None:
            return []

        candidates = self._find_table_candidate_nodes(tree.root)
        if not candidates:
            return []

        await self._log.info(
            f"[Compile] Targeted extraction: {len(candidates)} candidate "
            f"nodes in {Path(file_path).name}"
        )

        results: list[dict] = []
        seen_pages: set[int] = set()

        for node in candidates:
            if node.page_range is None:
                continue
            start_page, end_page = node.page_range
            # Skip pages already processed by another candidate
            page_nums = [p for p in range(start_page, end_page + 1)
                         if p not in seen_pages]
            if not page_nums:
                continue

            try:
                pages = DocumentExtractor.extract_page_range(
                    file_path, start_page, end_page,
                )
            except Exception as exc:
                await self._log.warning(
                    f"[Compile] Targeted extraction page read failed "
                    f"({start_page}-{end_page}): {exc}"
                )
                continue

            for pc in pages:
                if pc.page_number in seen_pages:
                    continue
                seen_pages.add(pc.page_number)
                # Numeric density gate – skip pages unlikely to contain tables
                if not self._page_has_table_density(pc.content):
                    continue
                regions = self._identify_table_regions(pc.content)
                for region in regions:
                    truncated = region[:_TARGETED_TABLE_MAX_CHARS]
                    results.append({
                        "page": pc.page_number,
                        "content": truncated,
                        "source": f"targeted:{node.title[:80]}",
                    })

        return results

    def _find_table_candidate_nodes(
        self, root: "TreeNode",
    ) -> list["TreeNode"]:
        """Collect leaf nodes that likely contain tables.

        Uses generic, domain-agnostic structural signals (any match
        suffices):

        - ``node.content_type == "table"`` – already tagged during compile.
        - ``node.table_count > 0`` – known to contain tables.
        - Has a valid ``page_range`` with span ≤ ``_TABLE_PAGE_SPAN_LIMIT``.
        """
        candidates: list = []

        def _walk(node: "TreeNode") -> None:
            if node.leaf:
                # Signal 1: content_type marked as table
                if getattr(node, "content_type", None) == "table":
                    candidates.append(node)
                    return
                # Signal 2: known to contain tables
                if getattr(node, "table_count", 0) > 0:
                    candidates.append(node)
                    return
                # Signal 3: moderate page span (tables rarely span many pages)
                page_range = getattr(node, "page_range", None)
                if page_range and len(page_range) == 2:
                    span = page_range[1] - page_range[0] + 1
                    if 1 <= span <= _TABLE_PAGE_SPAN_LIMIT:
                        candidates.append(node)
            else:
                for child in node.children:
                    _walk(child)

        _walk(root)
        return candidates

    # ------------------------------------------------------------------ #
    #  LLM-based heading normalisation                                     #
    # ------------------------------------------------------------------ #

    @staticmethod
    def _extract_heading_candidates(
        content: str,
    ) -> list[tuple[re.Match, str, str]]:
        """Extract candidate lines that *might* be section headings.

        Returns a list of ``(match, title_text, source_tag)`` triples
        where *source_tag* is ``"bold"`` or ``"standalone"``.

        Bold lines (``**Title**``) are always candidates.  Short
        standalone lines (surrounded by blank lines, 10-100 chars) are
        included only when they pass structural heuristics that filter
        out data rows, sentences, and existing headings.
        """
        occupied: list[tuple[int, int]] = []
        candidates: list[tuple[re.Match, str, str]] = []

        def _overlaps(start: int, end: int) -> bool:
            return any(s < end and start < e for s, e in occupied)

        for m in _BOLD_LINE_RE.finditer(content):
            title = m.group(1).strip()
            if title and not _overlaps(m.start(), m.end()):
                occupied.append((m.start(), m.end()))
                candidates.append((m, title, "bold"))

        for m in _STANDALONE_LINE_RE.finditer(content):
            text = m.group(1).strip()
            if len(text) < 10:
                continue
            text_offset = m.start() + m.group(0).index(m.group(1))
            if _overlaps(text_offset, text_offset + len(m.group(1))):
                continue
            if text.startswith(("#", "**")):
                continue
            if _NUM_TOKEN_RE.search(text):
                continue
            if text.endswith((".", "。", "!", "?", "！", "？")):
                continue
            if len(text.split()) > 12:
                continue
            occupied.append((text_offset, text_offset + len(m.group(1))))
            candidates.append((m, text, "standalone"))

        candidates.sort(key=lambda t: t[0].start())
        return candidates[:_HEADING_CANDIDATE_CAP]

    async def _normalize_bold_headings(self, content: str) -> str:
        """Detect and promote bold/standalone section titles to headings.

        Three-phase pipeline:
          1. **Extract** candidate lines via regex (deterministic).
          2. **Classify** candidates with a single LLM call — the LLM
             returns which indices are section headings and their level.
          3. **Replace** confirmed headings deterministically.

        Short-circuits when no candidates are found (zero LLM calls).
        On any LLM / parse failure, returns the original content unchanged
        (graceful degradation — equivalent to no-op).

        The transformation is idempotent: existing ``#`` headings never
        enter the candidate set.
        """
        if not content:
            return content

        candidates = self._extract_heading_candidates(content)
        if not candidates:
            return content

        listing = "\n".join(
            f"{i}: \"{title}\"" for i, (_, title, _tag) in enumerate(candidates)
        )

        from sirchmunk.llm.prompts import COMPILE_CLASSIFY_HEADINGS
        prompt = COMPILE_CLASSIFY_HEADINGS.format(candidates=listing)

        try:
            resp = await self._llm.achat(
                [{"role": "user", "content": prompt}],
            )
            raw = resp.content.strip()
            headings = self._parse_heading_classifications(raw, len(candidates))
        except Exception:
            return content

        if not headings:
            return content

        return self._apply_heading_promotions(content, candidates, headings)

    @staticmethod
    def _parse_heading_classifications(
        raw: str,
        num_candidates: int,
    ) -> list[tuple[int, int]]:
        """Parse LLM JSON response into a list of ``(idx, level)`` pairs.

        Robustly handles markdown code fences, trailing commas, and
        out-of-range indices.  Returns an empty list on any parse failure.
        """
        cleaned = raw.strip()
        if cleaned.startswith("```"):
            lines = cleaned.splitlines()
            lines = [ln for ln in lines if not ln.strip().startswith("```")]
            cleaned = "\n".join(lines).strip()

        try:
            items = json.loads(cleaned)
        except json.JSONDecodeError:
            m = re.search(r"\[.*\]", cleaned, re.DOTALL)
            if not m:
                return []
            try:
                items = json.loads(m.group())
            except json.JSONDecodeError:
                return []

        if not isinstance(items, list):
            return []

        result: list[tuple[int, int]] = []
        for item in items:
            if isinstance(item, dict):
                idx = item.get("idx")
                level = item.get("level", 2)
            elif isinstance(item, int):
                idx, level = item, 2
            else:
                continue
            if not isinstance(idx, int) or not (0 <= idx < num_candidates):
                continue
            level = max(2, min(4, int(level)))
            result.append((idx, level))
        return result

    @staticmethod
    def _apply_heading_promotions(
        content: str,
        candidates: list[tuple[re.Match, str, str]],
        headings: list[tuple[int, int]],
    ) -> str:
        """Apply heading promotions to *content* in reverse-offset order.

        Processes replacements from end-to-start so that earlier offsets
        remain valid after each substitution.
        """
        heading_map: dict[int, int] = dict(headings)

        replacements: list[tuple[int, int, str]] = []
        for idx, (match, title, tag) in enumerate(candidates):
            if idx not in heading_map:
                continue
            level = heading_map[idx]
            prefix = "#" * level
            if tag == "bold":
                replacements.append((match.start(), match.end(), f"{prefix} {title}"))
            else:
                text_start = match.start() + match.group(0).index(match.group(1))
                text_end = text_start + len(match.group(1))
                replacements.append((text_start, text_end, f"{prefix} {title}"))

        replacements.sort(key=lambda r: r[0], reverse=True)
        for start, end, replacement in replacements:
            content = content[:start] + replacement + content[end:]
        return content

    @staticmethod
    def _page_has_table_density(page_text: str) -> bool:
        """Return True if *page_text* likely contains tabular numeric data.

        Two independent signals (either suffices):

        1. **Character-level density** — fraction of digit/symbol chars
           relative to total non-whitespace exceeds the threshold.
        2. **Token-dense line** — any single line contains
           ``_DENSE_LINE_MIN_TOKENS`` or more numeric tokens, which
           catches pages where pypdf flattens all content into ≤ 2 lines.
        """
        if not page_text:
            return False
        non_ws = sum(1 for ch in page_text if not ch.isspace())
        if non_ws == 0:
            return False
        numeric_chars = sum(
            1 for ch in page_text
            if ch.isdigit() or ch in "$%(),.+-"
        )
        if (numeric_chars / non_ws) >= _TABLE_NUMERIC_DENSITY_THRESHOLD:
            return True
        return any(
            len(_NUM_TOKEN_RE.findall(line)) >= _DENSE_LINE_MIN_TOKENS
            for line in page_text.split("\n")
        )

    @staticmethod
    def _identify_table_regions(page_text: str) -> list[str]:
        """Identify contiguous table-like regions in *page_text*.

        Two complementary strategies:

        1. **Consecutive-line detection** — a run of ≥ 3 lines each
           containing ≥ 2 numeric tokens forms a table region.  Works
           well when pypdf preserves per-row line breaks.
        2. **Dense-line detection** — a *single* line with ≥
           ``_DENSE_LINE_MIN_TOKENS`` numeric tokens is treated as a
           table region.  This handles PDFs where pypdf collapses
           the entire page into one or two very long lines.

        Returns:
            List of extracted region strings (may be empty).
        """
        if not page_text:
            return []

        _MIN_NUMS_PER_LINE = 2
        _MIN_CONSECUTIVE = 3

        lines = page_text.split("\n")
        token_counts = [
            len(_NUM_TOKEN_RE.findall(line)) for line in lines
        ]

        regions: list[str] = []
        captured_lines: set[int] = set()

        # --- Strategy 1: consecutive-line runs ---
        run_start: int | None = None
        for i, cnt in enumerate(token_counts):
            if cnt >= _MIN_NUMS_PER_LINE:
                if run_start is None:
                    run_start = i
            else:
                if run_start is not None:
                    if i - run_start >= _MIN_CONSECUTIVE:
                        start = max(0, run_start - 1)
                        end = min(len(lines), i + 1)
                        regions.append(
                            "\n".join(lines[start:end]).strip()
                        )
                        captured_lines.update(range(start, end))
                    run_start = None
        if run_start is not None and len(lines) - run_start >= _MIN_CONSECUTIVE:
            start = max(0, run_start - 1)
            regions.append("\n".join(lines[start:]).strip())
            captured_lines.update(range(start, len(lines)))

        # --- Strategy 2: dense-line detection ---
        for i, cnt in enumerate(token_counts):
            if cnt >= _DENSE_LINE_MIN_TOKENS and i not in captured_lines:
                start = max(0, i - 1)
                end = min(len(lines), i + 2)
                regions.append("\n".join(lines[start:end]).strip())

        return regions

    @staticmethod
    def _get_table_page(entry: dict) -> int | None:
        """统一获取表格条目的页码，兼容 page_number 和 page 两种字段名。"""
        p = entry.get("page_number") or entry.get("page")
        return int(p) if p is not None else None

    @classmethod
    def _merge_table_digests(
        cls, existing: list[dict], new_tables: list[dict],
    ) -> list[dict]:
        """Merge *new_tables* into *existing* digest, deduplicating by page.

        If an existing entry and a new entry share the same page number,
        the new entry is skipped (existing kreuzberg-detected table takes
        precedence because it has richer structure like cells/markdown).

        Returns:
            Merged list suitable for storage in the table-digest JSON.
        """
        existing_pages = {cls._get_table_page(e) for e in existing}
        existing_pages.discard(None)

        merged = list(existing)
        for tbl in new_tables:
            page = cls._get_table_page(tbl)
            if page is not None and page in existing_pages:
                continue
            merged.append({
                "page_number": page,
                "markdown": tbl.get("markdown", "") or tbl.get("content", ""),
                "row_count": tbl.get("row_count"),
                "col_count": tbl.get("col_count"),
                "cells": tbl.get("cells", []),
                "source": tbl.get("source", "supplementary"),
            })
        return merged

    async def _supplement_table_digest(
        self,
        file_path: str,
        new_tables: list[dict],
        result: "FileCompileResult",
        *,
        source_label: str,
    ) -> None:
        """Merge supplementary tables into the persisted table digest.

        Loads the existing digest (if any), merges *new_tables* with
        page-level deduplication, and writes the updated digest back.
        Updates *result* metadata in place.
        """
        if not new_tables:
            return

        file_hash = get_fast_hash(file_path) or ""
        if not file_hash:
            return

        digest_dir = self._compile_dir / "table_digests"
        digest_path = digest_dir / f"{file_hash}.json"

        existing: list[dict] = []
        if result.has_table_digest and digest_path.exists():
            try:
                raw = json.loads(digest_path.read_text(encoding="utf-8"))
                existing = raw.get("tables", [])
            except Exception:
                pass

        merged = self._merge_table_digests(existing, new_tables)
        if not merged:
            return

        digest_dir.mkdir(parents=True, exist_ok=True)
        digest_path.write_text(
            json.dumps(
                {"version": 1, "table_count": len(merged), "tables": merged},
                ensure_ascii=False,
            ),
            encoding="utf-8",
        )
        result.has_table_digest = True
        result.table_count = len(merged)
        await self._log.info(
            f"[Compile] {source_label}: +{len(new_tables)} tables for "
            f"{Path(file_path).name} (total={len(merged)})"
        )

    def _get_covered_table_pages(self, file_path: str) -> Set[int]:
        """Return the set of page numbers already present in the table digest."""
        file_hash = get_fast_hash(file_path) or ""
        if not file_hash:
            return set()

        digest_path = (
            self._compile_dir / "table_digests" / f"{file_hash}.json"
        )
        if not digest_path.exists():
            return set()

        try:
            raw = json.loads(digest_path.read_text(encoding="utf-8"))
            pages: Set[int] = set()
            for t in raw.get("tables", []):
                p = self._get_table_page(t)
                if p is not None:
                    pages.add(p)
            return pages
        except Exception:
            return set()

    # ------------------------------------------------------------------ #
    #  P1: Enrich table digest with ENHANCED content                       #
    # ------------------------------------------------------------------ #

    @staticmethod
    def _build_page_char_map(
        tree_root: Any,
        max_page_span: int = _TABLE_PAGE_SPAN_LIMIT,
    ) -> Dict[int, Tuple[int, int]]:
        """Map page numbers to ``(start_char, end_char)`` in ENHANCED content.

        Aggregates ``char_range`` bounds from leaf nodes whose
        ``page_range`` intersects a given page.  To avoid inflated
        ranges from wide-spanning nodes (e.g. a cover-page node
        spanning pages 1–85), only nodes with a page span ≤
        *max_page_span* are used when available; wider nodes serve
        as a fallback.
        """
        # (char_start, char_end, page_span) per page
        entries: Dict[int, List[Tuple[int, int, int]]] = {}

        def _walk(node: Any) -> None:
            children = getattr(node, "children", None) or []
            if isinstance(node, dict):
                children = node.get("children", [])
            if not children:
                pr = (
                    getattr(node, "page_range", None)
                    if not isinstance(node, dict)
                    else node.get("page_range")
                )
                cr = (
                    getattr(node, "char_range", None)
                    if not isinstance(node, dict)
                    else node.get("char_range")
                )
                if (
                    pr
                    and cr
                    and len(pr) >= 2
                    and len(cr) >= 2
                ):
                    span = int(pr[1]) - int(pr[0]) + 1
                    for p in range(int(pr[0]), int(pr[1]) + 1):
                        entries.setdefault(p, []).append(
                            (int(cr[0]), int(cr[1]), span)
                        )
            for ch in children:
                _walk(ch)

        _walk(tree_root)

        result: Dict[int, Tuple[int, int]] = {}
        for page, elist in entries.items():
            narrow = [e for e in elist if e[2] <= max_page_span]
            chosen = narrow if narrow else elist
            result[page] = (
                min(e[0] for e in chosen),
                max(e[1] for e in chosen),
            )
        return result

    @staticmethod
    def _find_enhanced_region(
        enhanced_content: str,
        pypdf_text: str,
        budget: int = _TARGETED_TABLE_MAX_CHARS,
    ) -> Optional[str]:
        """Locate the ENHANCED content region matching *pypdf_text*.

        Uses progressively shorter text anchors extracted from the
        pypdf content to find the corresponding position in the
        ENHANCED (kreuzberg markdown) text.  Whitespace is normalised
        in the anchor to handle formatting differences (pypdf line
        breaks vs kreuzberg markdown spacing).  This avoids reliance
        on page-number alignment, which may differ between the two
        extractors.

        Returns the ENHANCED slice (up to *budget* chars) or ``None``.
        """
        text = pypdf_text.strip()
        for prefix in ("Table of Contents\n", "Table of Contents "):
            if text.startswith(prefix):
                text = text[len(prefix):]
        text = text.strip()

        for anchor_len in (80, 50, 30):
            raw = text[:anchor_len].strip()
            if len(raw) < 15:
                continue
            anchor = " ".join(raw.split())
            pos = enhanced_content.find(anchor)
            if pos < 0:
                continue
            start = max(
                0,
                enhanced_content.rfind("\n", max(0, pos - 300), pos) + 1,
            )
            end = min(len(enhanced_content), start + budget)
            return enhanced_content[start:end].strip()

        return None

    def _enrich_table_digest_content(
        self,
        file_path: str,
        enhanced_content: str,
        tree_root: Optional[Any],
    ) -> None:
        """Replace pypdf-sourced table text with ENHANCED content slices.

        Targeted extraction tables use pypdf, which often produces dense
        single-line text (the "2-line page" problem).  This method
        locates each table's content in the ENHANCED (kreuzberg markdown)
        text via anchor matching and replaces the ``markdown`` field when
        the ENHANCED version has substantially better structure.

        Only tables whose ``source`` indicates pypdf origin are
        candidates; kreuzberg-detected tables already have high-quality
        markdown and are left untouched.
        """
        if not enhanced_content:
            return

        file_hash = get_fast_hash(file_path) or ""
        if not file_hash:
            return

        digest_path = (
            self._compile_dir / "table_digests" / f"{file_hash}.json"
        )
        if not digest_path.exists():
            return

        try:
            raw = json.loads(digest_path.read_text(encoding="utf-8"))
            tables = raw.get("tables", [])
        except Exception:
            return

        if not tables:
            return

        modified = False
        for table in tables:
            source = table.get("source", "")
            if not (
                source.startswith("targeted:")
                or source == "content_scan"
            ):
                continue

            current = table.get("markdown", "")
            if not current:
                continue

            enhanced_region = self._find_enhanced_region(
                enhanced_content, current,
            )
            if not enhanced_region:
                continue

            current_lines = len(current.strip().split("\n"))
            enhanced_lines = len(enhanced_region.split("\n"))

            if enhanced_lines > max(current_lines, 3):
                table["markdown"] = enhanced_region[
                    :_TARGETED_TABLE_MAX_CHARS
                ]
                modified = True

        if modified:
            digest_path.write_text(
                json.dumps(raw, ensure_ascii=False),
                encoding="utf-8",
            )

    # ------------------------------------------------------------------ #
    #  Tree-independent content-based table scanning                       #
    # ------------------------------------------------------------------ #

    async def _content_based_table_scan(
        self,
        file_path: str,
        total_pages: Optional[int],
        covered_pages: Set[int],
        *,
        enhanced_content: Optional[str] = None,
        tree_root: Optional[Any] = None,
    ) -> list[dict]:
        """Scan PDF pages for table-like regions via numeric density.

        Uses a two-tier strategy:

        1. **pypdf page scan** — reads every page individually.  Works well
           when pypdf preserves per-row line breaks.
        2. **ENHANCED content fallback** — if pypdf yields poor line
           structure (> 50 % of pages have ≤ 3 lines), falls back to
           scanning the kreuzberg ENHANCED markdown content, which often
           has better formatting.  Page numbers are recovered via the
           tree's ``char_range → page_range`` mapping.

        Args:
            file_path:          Path to the PDF file.
            total_pages:        Total page count.
            covered_pages:      Page numbers already in the table digest.
            enhanced_content:   Cached kreuzberg ENHANCED text (optional).
            tree_root:          Tree root node for char → page mapping (optional).

        Returns:
            List of table dicts compatible with the digest format.
        """
        if not total_pages or total_pages <= 0:
            return []

        results = await self._pypdf_page_scan(
            file_path, total_pages, covered_pages,
        )

        if results or not enhanced_content or not tree_root:
            return results

        return self._enhanced_content_scan(
            enhanced_content, total_pages, covered_pages, tree_root,
        )

    async def _pypdf_page_scan(
        self,
        file_path: str,
        total_pages: int,
        covered_pages: Set[int],
    ) -> list[dict]:
        """Primary scan: per-page pypdf extraction with density heuristics.

        Pages are loaded in batches of ``_PAGE_SCAN_BATCH_SIZE`` to bound
        peak memory when processing large PDFs (200-400+ pages).
        """
        results: list[dict] = []
        poor_line_count = 0

        for batch_start in range(1, total_pages + 1, _PAGE_SCAN_BATCH_SIZE):
            batch_end = min(batch_start + _PAGE_SCAN_BATCH_SIZE, total_pages + 1)
            batch_pages = list(range(batch_start, batch_end))
            try:
                pages = DocumentExtractor.extract_pages(file_path, batch_pages)
            except Exception as exc:
                await self._log.warning(
                    f"[Compile] Content-based scan: page read failed for "
                    f"{Path(file_path).name}: {exc}"
                )
                return []

            for pc in pages:
                if len(pc.content.split("\n")) <= 3:
                    poor_line_count += 1
                if pc.page_number in covered_pages:
                    continue
                if not self._page_has_table_density(pc.content):
                    continue
                for region in self._identify_table_regions(pc.content):
                    results.append({
                        "page": pc.page_number,
                        "content": region[:_TARGETED_TABLE_MAX_CHARS],
                        "source": "content_scan",
                    })
            del pages

        if results:
            return results

        if poor_line_count > total_pages * 0.5:
            return []

        return results

    @staticmethod
    def _enhanced_content_scan(
        enhanced_content: str,
        total_pages: int,
        covered_pages: Set[int],
        tree_root: Any,
    ) -> list[dict]:
        """Fallback scan: use ENHANCED (kreuzberg markdown) content.

        Scans the full ENHANCED text line-by-line for dense-token lines,
        then maps each detected region back to a page number using the
        tree's ``char_range → page_range`` mapping.
        """
        char_page_map = KnowledgeCompiler._build_char_to_page_map(
            tree_root, total_pages,
        )
        if not char_page_map:
            return []

        breakpoints = [cp[0] for cp in char_page_map]

        results: list[dict] = []
        offset = 0
        for line in enhanced_content.split("\n"):
            token_count = len(_NUM_TOKEN_RE.findall(line))
            if token_count >= _DENSE_LINE_MIN_TOKENS:
                idx = bisect.bisect_right(breakpoints, offset) - 1
                page = char_page_map[max(0, idx)][1] if idx >= 0 else 1
                if page not in covered_pages:
                    results.append({
                        "page": page,
                        "content": line[:_TARGETED_TABLE_MAX_CHARS],
                        "source": "content_scan:enhanced",
                    })
                    covered_pages.add(page)
            offset += len(line) + 1  # +1 for '\n'

        return results

    @staticmethod
    def _build_char_to_page_map(
        tree_root: Any,
        total_pages: int,
    ) -> list[tuple[int, int]]:
        """Build a sorted (char_start, page_number) list from tree leaves.

        Enables efficient binary-search lookup from any character offset
        in the ENHANCED content to the corresponding page number.
        """
        entries: list[tuple[int, int]] = []

        def _collect(node: Any) -> None:
            children = getattr(node, "children", None) or []
            if isinstance(node, dict):
                children = node.get("children", [])
            pr = (
                getattr(node, "page_range", None)
                if not isinstance(node, dict)
                else node.get("page_range")
            )
            cr = (
                getattr(node, "char_range", None)
                if not isinstance(node, dict)
                else node.get("char_range")
            )
            if not children and cr and pr:
                page = pr[0] if isinstance(pr, (list, tuple)) else pr
                char_start = cr[0] if isinstance(cr, (list, tuple)) else cr
                if page and char_start is not None:
                    entries.append((int(char_start), int(page)))
            for ch in children:
                _collect(ch)

        _collect(tree_root)

        if not entries:
            return [(0, 1)]
        entries.sort()
        return entries

    def _find_force_ocr_candidates(
        self,
        file_path: str,
        total_pages: Optional[int],
        covered_pages: Set[int],
    ) -> List[int]:
        """Identify pages worth re-extracting with forced OCR.

        Returns 0-indexed page numbers for pages that have high numeric
        density (suggesting tabular content) but are NOT already covered
        by any table in the digest.  The result is capped at
        :data:`_FORCE_OCR_MAX_PAGES`.
        """
        if not total_pages or total_pages <= 0:
            return []

        all_page_nums = list(range(1, total_pages + 1))
        try:
            pages = DocumentExtractor.extract_pages(file_path, all_page_nums)
        except Exception:
            return []

        candidates: List[int] = []
        for pc in pages:
            if pc.page_number in covered_pages:
                continue
            if self._page_has_table_density(pc.content):
                candidates.append(pc.page_number - 1)  # 0-indexed for kreuzberg

        return sorted(candidates)[:_FORCE_OCR_MAX_PAGES]

    # ------------------------------------------------------------------ #
    #  Selective force-OCR re-extraction (P2)                              #
    # ------------------------------------------------------------------ #

    async def _selective_force_ocr_tables(
        self,
        file_path: str,
        gap_pages: List[int],
    ) -> list[dict[str, Any]]:
        """Extract text from gap pages using pypdf (no kreuzberg re-call).

        Earlier versions spawned a second kreuzberg extraction with
        ``force_ocr_pages``, which doubled native memory pressure.
        Using pypdf instead avoids Rust/native allocations entirely
        while still capturing page text for the table digest.

        Args:
            file_path:  Path to the PDF.
            gap_pages:  0-indexed page numbers.

        Returns:
            List of table-compatible dicts (``markdown``, ``page_number``).
        """
        if not gap_pages:
            return []

        capped = sorted(gap_pages)[:_FORCE_OCR_MAX_PAGES]
        one_indexed = [p + 1 for p in capped]
        try:
            pages = DocumentExtractor.extract_pages(file_path, one_indexed)
        except Exception:
            return []

        tables: list[dict[str, Any]] = []
        for pc in pages:
            text = (pc.content or "").strip()
            if text and self._page_has_table_density(text):
                tables.append({
                    "markdown": text,
                    "cells": [],
                    "page_number": pc.page_number,
                })
        return tables

    # ------------------------------------------------------------------ #
    #  Summary index for embedding + BM25 fallback                        #
    # ------------------------------------------------------------------ #

    async def _build_summary_index(self, manifest: CompileManifest) -> None:
        """Build summary embedding + BM25 index for fallback search.

        Creates a lightweight index mapping each compiled file to:
        - Its summary text
        - Pre-computed embedding vector (384-dim, if EmbeddingUtil available)
        - Tokenized summary with term frequencies (via TokenizerUtil)

        The index is saved to .cache/compile/summary_index.json and consumed
        by search.py as a last-resort fallback when rga keyword search fails.

        Reuses ``self._embedding`` when available to avoid loading a duplicate
        model into memory.  Falls back to a fresh instance otherwise.
        """
        try:
            from sirchmunk.utils.tokenizer_util import TokenizerUtil
            from sirchmunk.learnings.summary_index import CompileSummaryIndex, SummaryIndexEntry

            entries: List[SummaryIndexEntry] = []
            summaries: List[str] = []

            for file_path, entry in manifest.files.items():
                if entry.summary:
                    entries.append(SummaryIndexEntry(
                        file_path=file_path,
                        summary=entry.summary,
                    ))
                    summaries.append(entry.summary)

            if not entries:
                return

            tokenizer = TokenizerUtil()
            for idx, entry in enumerate(entries):
                tokens = tokenizer.segment(entry.summary)
                entry.tokens = tokens
                entry.token_freqs = {}
                for t in tokens:
                    entry.token_freqs[t] = entry.token_freqs.get(t, 0) + 1

            # Reuse the compiler's embedding client to avoid duplicate model load
            try:
                embedding_util = self._embedding
                if embedding_util is None:
                    from sirchmunk.utils.embedding_util import EmbeddingUtil
                    embedding_util = EmbeddingUtil()
                    embedding_util.start_loading()

                await embedding_util._ensure_model_async(timeout=60)

                if embedding_util.is_ready():
                    embeddings = await embedding_util.embed(summaries)
                    for i, emb in enumerate(embeddings):
                        entries[i].embedding = emb
                    await self._log.info(
                        f"Summary index: computed embeddings for {len(entries)} entries"
                    )
            except Exception as emb_exc:
                await self._log.warning(
                    f"Summary index: embedding computation skipped: {emb_exc}"
                )

            index = CompileSummaryIndex(entries)
            index.save(self._compile_dir / "summary_index.json")

        except Exception as exc:
            await self._log.warning(f"Failed to build summary index: {exc}")

    # ------------------------------------------------------------------ #
    #  Manifest I/O                                                       #
    # ------------------------------------------------------------------ #

    def _load_manifest(self) -> CompileManifest:
        if self._manifest_path.exists():
            try:
                return CompileManifest.from_json(
                    self._manifest_path.read_text(encoding="utf-8")
                )
            except Exception:
                pass
        return CompileManifest()

    def _save_manifest(self, manifest: CompileManifest) -> None:
        """Atomically persist the manifest via write-to-tmp + rename.

        This prevents partial JSON on disk if the process is killed mid-write.
        """
        tmp_path = self._manifest_path.with_suffix(".json.tmp")
        tmp_path.write_text(manifest.to_json(), encoding="utf-8")
        tmp_path.replace(self._manifest_path)

    # ------------------------------------------------------------------ #
    #  Document catalog for search-time routing                           #
    # ------------------------------------------------------------------ #

    def _build_document_catalog(self, manifest: CompileManifest) -> None:
        """Generate a lightweight catalog mapping files to their tree root summaries.

        The catalog is consumed by FAST search to fuse query analysis with
        LLM-driven document routing in a single prompt.  Each entry carries
        the filename and a truncated root summary (<= _MANIFEST_SUMMARY_MAX_LEN chars).

        Summary is sourced from the manifest (populated during Phase 2 compile),
        with a tree-root fallback for backward compatibility.
        """
        tree_cache = self._compile_dir / "trees"
        entries: List[Dict[str, str]] = []

        for file_path, entry in manifest.files.items():
            summary = entry.summary  # Primary: manifest-persisted summary

            # Fallback: read from tree root if manifest summary is empty
            if not summary and entry.has_tree and tree_cache.exists():
                tree_file = tree_cache / f"{entry.file_hash}.json"
                if tree_file.exists():
                    try:
                        tree = DocumentTree.from_json(
                            tree_file.read_text(encoding="utf-8"),
                        )
                        if tree.root and tree.root.summary:
                            summary = tree.root.summary[:_MANIFEST_SUMMARY_MAX_LEN]
                    except Exception:
                        pass

            entries.append({
                "path": file_path,
                "name": Path(file_path).name,
                "summary": summary,
            })

        catalog_path = self._compile_dir / "document_catalog.json"
        catalog_path.write_text(
            json.dumps(entries, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
